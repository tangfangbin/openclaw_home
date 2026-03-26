#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
柳州城市介绍 PowerPoint 演示文稿 - 每页配图版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(30, 100, 160)
DARK_BLUE = RGBColor(20, 60, 120)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)
ORANGE = RGBColor(255, 140, 0)

IMG_DIR = "/Users/blade/.openclaw/workspace/liuzhou_ppt"

def add_cover_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background image
    if os.path.exists(f"{IMG_DIR}/p1_cover.jpg"):
        slide.shapes.add_picture(f"{IMG_DIR}/p1_cover.jpg", 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(3))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(20, 60, 120)
    overlay.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_image_content_slide(prs, page_num, title, image_path, content_list, emoji="📍"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    
    # Header background
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(1.3))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(240, 245, 250)
    header_bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Page number
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.35), Inches(0.8), Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 200, 220)
    p.alignment = PP_ALIGN.RIGHT
    
    # Image (left side)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.7), width=Inches(5.5))
    
    # Content (right side)
    content_box = slide.shapes.add_textbox(Inches(6.3), Inches(2), Inches(6.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(14)
    
    return slide

def add_two_column_slide(prs, page_num, title, image_path, left_content, right_content, emoji="🏛️"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    
    # Header background
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(1.3))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(240, 245, 250)
    header_bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Page number
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.35), Inches(0.8), Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 200, 220)
    p.alignment = PP_ALIGN.RIGHT
    
    # Image (top center)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(3.5), Inches(1.6), width=Inches(6))
    
    # Left column content
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.6))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = left_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(5.5), Inches(1.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
    
    # Right column content
    right_title = slide.shapes.add_textbox(Inches(7), Inches(5.2), Inches(5.5), Inches(0.6))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = right_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    right_box = slide.shapes.add_textbox(Inches(7), Inches(5.8), Inches(5.5), Inches(1.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
    
    return slide

def add_ending_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background image
    if os.path.exists(f"{IMG_DIR}/p11_ending.jpg"):
        slide.shapes.add_picture(f"{IMG_DIR}/p11_ending.jpg", 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(20, 60, 120)
    overlay.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 创建幻灯片 ============

# 第1页：封面
add_cover_slide(prs, "柳州城市介绍", "桂中明珠·龙城柳州")

# 第2页：地理位置
add_image_content_slide(prs, 2, "地理位置",
    f"{IMG_DIR}/p2_location.jpg",
    [
        "位于广西壮族自治区中部偏东北",
        "地处桂中盆地，珠江水系柳江上游",
        "总面积 18,677 平方公里",
        "气候属于亚热带季风气候，温润多雨",
        "柳江呈"U"字形绕城而过"
    ], "📍")

# 第3页：行政区划
add_two_column_slide(prs, 3, "人口与行政区划",
    f"{IMG_DIR}/p3_admin.jpg",
    ["📊 人口数据", "常住人口约 415 万", "户籍人口约 393 万", "城镇化率约 65%"],
    ["🏛️ 行政区划", "5个区：城中、鱼峰、柳南、柳北、柳江", "5个县：柳城、鹿寨、融安、融水、三江"],
    "🏛️")

# 第4页：历史沿革
add_image_content_slide(prs, 4, "历史沿革",
    f"{IMG_DIR}/p4_history.jpg",
    [
        "公元前 111 年（西汉）始建城，已有 2100 多年历史",
        "因城南有一座柳江环绕的小山，取名'龙城'",
        "唐宋时期为柳州刺史驻地，商业渐兴",
        "明清时期为广西四大城市之一",
        "1949 年后成为广西重要工业基地"
    ], "📜")

# 第5页：支柱产业
add_two_column_slide(prs, 5, "支柱产业",
    f"{IMG_DIR}/p5_industry.jpg",
    ["🏭 重工业", "汽车工业（全国五大汽车城之一）", "钢铁工业（柳钢集团）", "工程机械（柳工集团）"],
    ["🧵 特色产业", "螺蛳粉产业（年产值超 500 亿）", "日化产业（两面针、田七）", "新能源产业"],
    "🏭")

# 第6页：风景名胜
add_image_content_slide(prs, 6, "风景名胜",
    f"{IMG_DIR}/p6_scenic.jpg",
    [
        "程阳八寨：世界文化遗产侗族村寨",
        "龙潭公园：集喀斯特山水于一体的大型公园",
        "柳侯公园：为纪念唐代文学家柳宗元而建",
        "百里柳江：夜景璀璨，两岸景观丰富",
        "融安石门仙湖：三江源头，生态完好"
    ], "🏞️")

# 第7页：民俗文化
add_image_content_slide(prs, 7, "民俗文化",
    f"{IMG_DIR}/p7_culture.jpg",
    [
        "侗族大歌：世界非物质文化遗产",
        "苗族芦笙舞：每年芦笙节热闹非凡",
        "壮族山歌：以歌会友，歌圩遍布",
        "三江程阳风雨桥：木构桥梁建筑奇迹",
        "非遗技艺：侗绣、苗锦、壮锦"
    ], "🎭")

# 第8页：特色美食
add_image_content_slide(prs, 8, "特色美食",
    f"{IMG_DIR}/p8_food.jpg",
    [
        "螺蛳粉：柳州第一美食，酸辣鲜香",
        "牛腊巴：柳城特产，香脆可口",
        "云片糕：甜而不腻，入口即化",
        "柳州酸：各式酸品小吃，开胃解暑",
        "露水汤圆：市区独有，咸口汤圆"
    ], "🍜")

# 第9页：交通优势
add_two_column_slide(prs, 9, "交通优势",
    f"{IMG_DIR}/p9_transport.jpg",
    ["🚄 铁路", "湘桂铁路、黔桂铁路交汇", "高铁可达南宁、桂林、广州", "规划建设柳肇铁路"],
    ["🚗 公路", "泉南高速、柳南高速交汇", "国道 209、322 穿城", "农村公路网络完善"],
    "🚄")

# 第10页：城市荣誉
add_image_content_slide(prs, 10, "城市荣誉",
    f"{IMG_DIR}/p10_honor.jpg",
    [
        "中国十大宜居城市",
        "国家历史文化名城",
        "全国双拥模范城",
        "国家园林城市",
        "广西高等教育基地"
    ], "🏆")

# 第11页：结束页
add_ending_slide(prs, "谢谢观看", "柳州——一座让人流连忘返的城市 🚀")

# 保存
output_path = "/Users/blade/Desktop/柳州城市介绍.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"文件大小: {os.path.getsize(output_path) / 1024:.1f} KB")
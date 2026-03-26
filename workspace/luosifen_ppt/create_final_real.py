#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
柳州螺蛳粉 PowerPoint 演示文稿 - 真实图片版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(200, 30, 30)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)

IMG_DIR = "/Users/blade/.openclaw/workspace/luosifen_ppt"

def add_title_slide(prs, title, subtitle, image_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(3))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(180, 30, 30)
    overlay.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_image_content_slide(prs, page_num, title, image_path, content_list, emoji="📍"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(1.3))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(250, 245, 240)
    header_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.35), Inches(0.8), Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 180, 180)
    p.alignment = PP_ALIGN.RIGHT
    
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.7), width=Inches(5.5))
    
    content_box = slide.shapes.add_textbox(Inches(6.3), Inches(2), Inches(6.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(14)
    
    return slide

def add_two_column_slide(prs, page_num, title, image_path, left_content, right_content, emoji="🏛️"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(1.3))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(250, 245, 240)
    header_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.35), Inches(0.8), Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 180, 180)
    p.alignment = PP_ALIGN.RIGHT
    
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(3.5), Inches(1.6), width=Inches(6))
    
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.6))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = left_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(5.5), Inches(1.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
    
    right_title = slide.shapes.add_textbox(Inches(7), Inches(5.2), Inches(5.5), Inches(0.6))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = right_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    right_box = slide.shapes.add_textbox(Inches(7), Inches(5.8), Inches(5.5), Inches(1.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
    
    return slide

def add_ending_slide(prs, title, subtitle, image_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(180, 30, 30)
    overlay.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 创建幻灯片 ============

# 第1页：封面
add_title_slide(prs, "柳州螺蛳粉", "一碗来自广西柳州的传奇风味", f"{IMG_DIR}/luosifen_food_1.jpg")

# 第2页：什么是螺蛳粉
add_image_content_slide(prs, 2, "什么是螺蛳粉？",
    f"{IMG_DIR}/luosifen_food_2.jpg",
    ["螺蛳粉是广西柳州的特色小吃", "以独特的酸、辣、鲜、烫口味著称", "主要原料：米粉、酸笋、花生、木耳", "灵魂汤底：用螺蛳、猪骨、香料熬制", "最大的特色是酸笋发酵的独特香味"], "🍜")

# 第3页：历史起源
add_image_content_slide(prs, 3, "历史起源",
    f"{IMG_DIR}/luosifen_food_3.jpg",
    ["起源时间：约20世纪80年代", "发源地：广西柳州市鱼峰区、胜利小区", "最早是路边摊小吃，因味道独特逐渐流行", "2014年，首家预包装螺蛳粉企业成立", "2020年成为'网红顶流'，走向世界"], "📜")

# 第4页：主要原料
add_two_column_slide(prs, 4, "主要原料",
    f"{IMG_DIR}/luosifen_food_4.jpg",
    ["🥣 汤底原料", "螺蛳（石螺）", "猪筒骨、鸡架", "八角、桂皮、草果", "香叶、小茴香"],
    ["🍜 配菜原料", "干米粉（柳州圆头粉）", "酸笋（发酵笋）", "腐竹，花生", "木耳、黄花菜"], "🥗")

# 第5页：制作方法
add_image_content_slide(prs, 5, "制作方法",
    f"{IMG_DIR}/luosifen_food_5.jpg",
    ["① 熬汤：螺蛳和猪骨一起熬制4-6小时", "② 泡粉：干米粉提前用温水泡软", "③ 焯水：米粉焯水后装碗", "④ 浇汤：浇上滚烫的螺蛳汤底", "⑤ 加配菜：依次放入酸笋、腐竹、花生等", "⑥ 调味：加入辣椒油、醋、蒜蓉等"], "👨‍🍳")

# 第6页：口味特点
add_two_column_slide(prs, 6, "口味特点",
    f"{IMG_DIR}/luosifen_food_6.jpg",
    ["🌶️ 辣", "以辣著称", "辣度可选微/中/重辣", "辣椒油是灵魂"],
    ["😋 鲜", "螺蛳汤底鲜甜", "猪骨胶原蛋白", "完美平衡辣味"], "⭐")

# 第7页：衍生吃法
add_image_content_slide(prs, 7, "衍生吃法与创新",
    f"{IMG_DIR}/luosifen_food_7.jpg",
    ["🥡 干捞螺蛳粉：不用汤底，拌着吃，更加入味", "🔥 炒螺蛳粉：用猛火炒制，口感更香", "🧊 凉拌螺蛳粉：夏季吃法，冰凉解暑", "🍳 螺蛳粉火锅：螺蛳粉汤底涮各种食材", "🥧 螺蛳粉月饼/粽子：黑暗料理界的网红"], "🎉")

# 第8页：柳州必吃推荐
add_image_content_slide(prs, 8, "柳州必吃推荐",
    f"{IMG_DIR}/luosifen_food_8.jpg",
    ["📍 胜利小区：螺蛳粉发源地之一，最正宗老味道", "📍 柳北区：多家网红老店聚集", "📍 窑埠古镇：新兴美食街区，环境好", "🏆 品牌推荐：好欢螺、螺霸王、柳江人家", "💡 建议：到柳州一定要去街边小店吃现煮的！"], "🗺️")

# 第9页：结束页
add_ending_slide(prs, "谢谢观看！", "柳州螺蛳粉，一口入魂 🚀", f"{IMG_DIR}/luosifen_food_10.jpg")

# 保存
output_path = "/Users/blade/Desktop/柳州螺蛳粉介绍_真图版.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"文件大小: {os.path.getsize(output_path) / 1024 / 1024:.1f} MB")
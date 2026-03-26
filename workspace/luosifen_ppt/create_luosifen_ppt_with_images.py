#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
柳州螺蛳粉 PowerPoint 演示文稿生成脚本 - 带图片版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
RED = RGBColor(200, 30, 30)
ORANGE = RGBColor(255, 120, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# 图片路径
IMG_DIR = "/Users/blade/.openclaw/workspace/luosifen_ppt"

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 红色背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    # 装饰圆
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(255, 140, 0)
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 100, 50)
    circle2.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, content_list, emoji="📍"):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(16)
    
    return slide

def add_image_slide(prs, title, image_path, description_list, emoji="🍜"):
    """添加图片+文字说明页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)
    
    # 图片
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.5), width=Inches(6))
    
    # 右侧说明
    content_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(description_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, emoji="🍜"):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)
    
    # 左栏标题
    left_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = left_content[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
    
    # 右栏标题
    right_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.6))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = right_content[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
    
    return slide

def add_closing_slide(prs, title, subtitle):
    """添加结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 全屏背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    # 装饰圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 140, 0)
    circle.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 开始创建幻灯片 ============

# 第1页：标题页
add_title_slide(prs, "柳州螺蛳粉", "一碗来自广西柳州的传奇风味")

# 第2页：什么是螺蛳粉（图文）
add_image_slide(prs, "什么是螺蛳粉？",
    f"{IMG_DIR}/luosifen_bowl.jpg",
    [
        "螺蛳粉是广西柳州的特色小吃",
        "以独特的酸、辣、鲜、烫口味著称",
        "主要原料：米粉、酸笋、花生、木耳",
        "灵魂汤底：用螺蛳、猪骨、香料熬制",
        "最大的特色是酸笋发酵的独特香味"
    ], "🍜")

# 第3页：历史起源
add_content_slide(prs, "历史起源", [
    "起源时间：约20世纪80年代",
    "发源地：广西柳州市鱼峰区、胜利小区",
    "最早是路边摊小吃，因味道独特逐渐流行",
    "2014年，首家预包装螺蛳粉企业成立",
    "2020年成为'网红顶流'，走向世界"
], "📜")

# 第4页：主要原料（双栏）
add_two_column_slide(prs, "主要原料",
    ["🥣 汤底原料", "螺蛳（石螺）", "猪筒骨、鸡架", "八角、桂皮、草果", "香叶、小茴香", "辣椒油"],
    ["🍜 配菜原料", "干米粉（柳州圆头粉）", "酸笋（发酵笋）", "腐竹、花生", "木耳、黄花菜", "萝卜干、酸豆角"],
    "🥗")

# 第5页：制作方法
add_content_slide(prs, "制作方法", [
    "① 熬汤：螺蛳和猪骨一起熬制4-6小时",
    "② 泡粉：干米粉提前用温水泡软",
    "③ 焯水：米粉焯水后装碗",
    "④ 浇汤：浇上滚烫的螺蛳汤底",
    "⑤ 加配菜：依次放入酸笋、腐竹、花生等",
    "⑥ 调味：加入辣椒油、醋、蒜蓉等"
], "👨‍🍳")

# 第6页：口味特点
add_two_column_slide(prs, "口味特点",
    ["🌶️ 辣", "柳州螺蛳粉以辣著称", "辣度可选择微辣、中辣、重辣", "辣椒油是灵魂调料之一"],
    ["😋 鲜", "螺蛳汤底鲜甜浓郁", "猪骨熬出胶原蛋白", "完美平衡辣味"],
    "⭐")

# 第7页：衍生吃法
add_content_slide(prs, "衍生吃法与创新", [
    "🥡 干捞螺蛳粉：不用汤底，拌着吃",
    "🔥 炒螺蛳粉：用猛火炒制，口感更香",
    "🧊 凉拌螺蛳粉：夏季吃法，冰凉解暑",
    "🍳 螺蛳粉火锅：螺蛳粉汤底涮各种食材",
    "🥧 螺蛳粉月饼/粽子：黑暗料理界的网红"
], "🎉")

# 第8页：美食地图
add_content_slide(prs, "柳州必吃推荐", [
    "📍 胜利小区：螺蛳粉发源地之一",
    "📍 柳北区：多家网红老店聚集",
    "📍 窑埠古镇：新兴美食街区",
    "🏆 重要品牌：好欢螺、螺霸王、柳江人家",
    "💡 建议：到柳州一定要去街边小店！"
], "🗺️")

# 第9页：结束页
add_closing_slide(prs, "谢谢观看！", "柳州螺蛳粉，一口入魂 🚀")

# 保存文件
output_path = "/Users/blade/Desktop/柳州螺蛳粉介绍_带图版.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
柳州螺蛳粉 PowerPoint 演示文稿 - 每页配图版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(220, 53, 69)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)
LIGHT_GRAY = RGBColor(245, 245, 245)

IMG_DIR = "/Users/blade/.openclaw/workspace/luosifen_ppt"

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()
    
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(255, 140, 0)
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 100, 50)
    circle2.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_with_image_slide(prs, title, image_path, content_list, emoji="📍"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # 图片（放在左侧或右侧，根据内容调整）
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.5), width=Inches(4.5))
    
    # 内容（右侧）
    content_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(7.3), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(14)
    
    return slide

def add_two_column_with_image_slide(prs, title, image_path, left_content, right_content, emoji="🍜"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # 图片（左侧中部）
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.5), width=Inches(3.8))
    
    # 左栏标题
    left_title = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(3.8), Inches(0.6))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = left_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(5), Inches(2.2), Inches(3.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title = slide.shapes.add_textbox(Inches(9.2), Inches(1.5), Inches(3.8), Inches(0.6))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = right_content[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(9.2), Inches(2.2), Inches(3.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content[1:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    
    return slide

def add_image_only_slide(prs, title, image_path, description, emoji="🍜"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # 大图（居中）
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(2), Inches(1.5), width=Inches(9))
    
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(1.2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_closing_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 220, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 创建幻灯片 ============

# 第1页：封面
add_title_slide(prs, "柳州螺蛳粉", "一碗来自广西柳州的传奇风味")

# 第2页：什么是螺蛳粉（配图）
add_content_with_image_slide(prs, "什么是螺蛳粉？",
    f"{IMG_DIR}/img_p2_what.jpg",
    [
        "螺蛳粉是广西柳州的特色小吃",
        "以独特的酸、辣、鲜、烫口味著称",
        "主要原料：米粉、酸笋、花生、木耳",
        "灵魂汤底：用螺蛳、猪骨、香料熬制",
        "最大的特色是酸笋发酵的独特香味",
        "2020年成为'网红顶流'，走向世界"
    ], "🍜")

# 第3页：历史起源（配图）
add_content_with_image_slide(prs, "历史起源",
    f"{IMG_DIR}/img_p3_history.jpg",
    [
        "起源时间：约20世纪80年代",
        "发源地：广西柳州市鱼峰区、胜利小区",
        "最早是路边摊小吃，因味道独特逐渐流行",
        "2014年，首家预包装螺蛳粉企业成立",
        "2020年成为'网红顶流'，年产值破百亿"
    ], "📜")

# 第4页：主要原料（配图+双栏）
add_two_column_with_image_slide(prs, "主要原料",
    f"{IMG_DIR}/img_p4_ingredients.jpg",
    ["🥣 汤底原料", "螺蛳（石螺）", "猪筒骨、鸡架", "八角、桂皮、草果", "香叶、小茴香"],
    ["🍜 配菜原料", "干米粉（柳州圆头粉）", "酸笋（发酵笋）", "腐竹、花生", "木耳、黄花菜"],
    "🥗")

# 第5页：制作方法（配图）
add_content_with_image_slide(prs, "制作方法",
    f"{IMG_DIR}/img_p5_cooking.jpg",
    [
        "① 熬汤：螺蛳和猪骨一起熬制4-6小时",
        "② 泡粉：干米粉提前用温水泡软",
        "③ 焯水：米粉焯水后装碗",
        "④ 浇汤：浇上滚烫的螺蛳汤底",
        "⑤ 加配菜：依次放入酸笋、腐竹、花生等",
        "⑥ 调味：加入辣椒油、醋、蒜蓉等"
    ], "👨‍🍳")

# 第6页：口味特点（配图+双栏）
add_two_column_with_image_slide(prs, "口味特点",
    f"{IMG_DIR}/img_p6_taste.jpg",
    ["🌶️ 辣", "以辣著称", "辣度可选微/中/重辣", "辣椒油是灵魂"],
    ["😋 鲜", "螺蛳汤底鲜甜", "猪骨胶原蛋白", "完美平衡辣味"],
    "⭐")

# 第7页：衍生吃法（配图）
add_content_with_image_slide(prs, "衍生吃法与创新",
    f"{IMG_DIR}/img_p7_innovations.jpg",
    [
        "🥡 干捞螺蛳粉：不用汤底，拌着吃，更加入味",
        "🔥 炒螺蛳粉：用猛火炒制，口感更香",
        "🧊 凉拌螺蛳粉：夏季吃法，冰凉解暑",
        "🍳 螺蛳粉火锅：螺蛳粉汤底涮各种食材",
        "🥧 螺蛳粉月饼/粽子：黑暗料理界的网红"
    ], "🎉")

# 第8页：柳州必吃推荐（配图）
add_content_with_image_slide(prs, "柳州必吃推荐",
    f"{IMG_DIR}/img_p8_map.jpg",
    [
        "📍 胜利小区：螺蛳粉发源地之一，最正宗老味道",
        "📍 柳北区：多家网红老店聚集",
        "📍 窑埠古镇：新兴美食街区，环境好",
        "🏆 品牌推荐：好欢螺、螺霸王、柳江人家",
        "💡 建议：到柳州一定要去街边小店吃现煮的！"
    ], "🗺️")

# 第9页：结束页
add_closing_slide(prs, "谢谢观看！", "柳州螺蛳粉，一口入魂 🚀")

# 保存
output_path = "/Users/blade/Desktop/柳州螺蛳粉介绍_每页配图版.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")